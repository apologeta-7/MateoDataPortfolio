"""
Build portfolio-facing artifacts from saved model evaluation outputs.

Generated outputs:
- dashboard/data_exports/*.csv
- presentation/assets/charts/*.png
- presentation/store_sales_forecasting_presentation.pptx
- presentation/store_sales_forecasting_presentation.pdf
"""

from __future__ import annotations

import shutil
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import Image, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle


ROOT = Path(__file__).resolve().parents[1]
TABLES = ROOT / "reports" / "tables"
FIGURES = ROOT / "reports" / "figures"
INFERENCE = ROOT / "reports" / "inference"
DASHBOARD_EXPORTS = ROOT / "dashboard" / "data_exports"
PRESENTATION_DIR = ROOT / "presentation"
PRESENTATION_CHARTS = PRESENTATION_DIR / "assets" / "charts"


TITLE = "Store Sales Forecasting"
SUBTITLE = "Corporacion Favorita demand forecast portfolio project"


def read_csv(path: Path) -> pd.DataFrame:
    if not path.exists():
        raise FileNotFoundError(path.as_posix())
    return pd.read_csv(path)


def round_numeric(dataframe: pd.DataFrame, digits: int = 6) -> pd.DataFrame:
    result = dataframe.copy()
    numeric_columns = result.select_dtypes(include="number").columns
    result[numeric_columns] = result[numeric_columns].round(digits)
    return result


def write_export(dataframe: pd.DataFrame, name: str) -> Path:
    DASHBOARD_EXPORTS.mkdir(parents=True, exist_ok=True)
    path = DASHBOARD_EXPORTS / name
    round_numeric(dataframe).to_csv(path, index=False)
    return path


def build_dashboard_exports() -> list[Path]:
    leaderboard = read_csv(TABLES / "final_model_leaderboard.csv")
    recommendation = read_csv(TABLES / "final_model_recommendation.csv")
    store = read_csv(TABLES / "final_error_by_store.csv")
    family = read_csv(TABLES / "final_error_by_family.csv")
    daily = read_csv(TABLES / "final_error_by_date.csv")
    promotion = read_csv(TABLES / "final_error_by_promotion.csv")
    holiday = read_csv(TABLES / "final_error_by_holiday.csv")
    feature_importance = read_csv(TABLES / "lightgbm_v1_feature_importance.csv")
    future_daily = read_csv(INFERENCE / "final_model_inference_v1_src_daily_summary.csv")

    lightgbm = leaderboard[leaderboard["model"] == "lightgbm_v1"].iloc[0]
    baseline = leaderboard[
        leaderboard["model"] == "baseline_rolling_mean_28_origin"
    ].iloc[0]
    prophet = leaderboard[leaderboard["model"] == "prophet_total_sales_v1"].iloc[0]
    daily_lightgbm = leaderboard[
        leaderboard["model"] == "lightgbm_v1_daily_aggregate"
    ].iloc[0]

    executive_kpis = pd.DataFrame(
        [
            {
                "recommended_model": recommendation.iloc[0]["recommended_model"],
                "forecast_horizon_days": int(
                    recommendation.iloc[0]["forecast_horizon_days"]
                ),
                "validation_start_date": lightgbm["validation_start_date"],
                "validation_end_date": lightgbm["validation_end_date"],
                "validation_rows": int(lightgbm["rows"]),
                "lightgbm_rmsle": lightgbm["rmsle"],
                "lightgbm_wape": lightgbm["wape"],
                "lightgbm_mae": lightgbm["mae"],
                "lightgbm_total_bias_pct": lightgbm["total_bias_pct"],
                "baseline_rmsle": baseline["rmsle"],
                "baseline_wape": baseline["wape"],
                "daily_lightgbm_rmsle": daily_lightgbm["rmsle"],
                "prophet_daily_rmsle": prophet["rmsle"],
                "future_forecast_rows": int(future_daily["rows"].sum()),
                "future_predicted_sales_sum": future_daily[
                    "predicted_sales_sum"
                ].sum(),
            }
        ]
    )

    store_columns = [
        "store_nbr",
        "city",
        "state",
        "store_type",
        "sales_volume_segment",
        "actual_total_sales_lightgbm",
        "predicted_total_sales_lightgbm",
        "rmsle_lightgbm",
        "wape_lightgbm",
        "total_bias_pct_lightgbm",
        "error_contribution_pct_lightgbm",
        "better_model_by_store_rmsle",
    ]
    family_columns = [
        "family",
        "sales_volume_segment",
        "zero_sales_pct",
        "promotion_rows_pct",
        "actual_total_sales_lightgbm",
        "predicted_total_sales_lightgbm",
        "rmsle_lightgbm",
        "wape_lightgbm",
        "total_bias_pct_lightgbm",
        "error_contribution_pct_lightgbm",
        "better_model_by_family_rmsle",
    ]

    segment_table = pd.concat(
        [
            promotion.assign(segment_type="promotion").rename(
                columns={"promotion_status": "segment"}
            ),
            holiday.assign(segment_type="holiday").rename(
                columns={"holiday_status": "segment"}
            ),
        ],
        ignore_index=True,
    )

    exports = [
        write_export(executive_kpis, "executive_kpis.csv"),
        write_export(leaderboard, "model_leaderboard.csv"),
        write_export(store[store_columns], "error_by_store.csv"),
        write_export(
            store[store_columns]
            .sort_values("error_contribution_pct_lightgbm", ascending=False)
            .head(15),
            "top_risk_stores.csv",
        ),
        write_export(family[family_columns], "error_by_family.csv"),
        write_export(
            family[family_columns]
            .sort_values("error_contribution_pct_lightgbm", ascending=False)
            .head(15),
            "top_risk_families.csv",
        ),
        write_export(daily, "error_by_date.csv"),
        write_export(segment_table, "promotion_holiday_segments.csv"),
        write_export(
            feature_importance.sort_values("importance_gain", ascending=False).head(20),
            "feature_importance_top20.csv",
        ),
        write_export(future_daily, "future_daily_forecast.csv"),
    ]

    readme = DASHBOARD_EXPORTS / "README.md"
    readme.write_text(
        "\n".join(
            [
                "# Dashboard Data Exports",
                "",
                "These CSV files are generated from the saved model evaluation outputs.",
                "They are intended as the Power BI import layer for the portfolio dashboard.",
                "",
                "Recommended Power BI pages:",
                "- Executive overview: `executive_kpis.csv`, `model_leaderboard.csv`.",
                "- Forecast horizon: `future_daily_forecast.csv`.",
                "- Store risk: `error_by_store.csv`, `top_risk_stores.csv`.",
                "- Family risk: `error_by_family.csv`, `top_risk_families.csv`.",
                "- Segment diagnostics: `promotion_holiday_segments.csv`.",
                "- Explainability: `feature_importance_top20.csv`.",
                "",
                "Regenerate with:",
                "",
                "```powershell",
                "..\\.venv\\Scripts\\python.exe scripts\\build_portfolio_artifacts.py",
                "```",
                "",
            ]
        ),
        encoding="utf-8",
    )
    exports.append(readme)
    return exports


def metric_lookup() -> dict[str, float | str | int]:
    leaderboard = read_csv(TABLES / "final_model_leaderboard.csv")
    recommendation = read_csv(TABLES / "final_model_recommendation.csv")
    future_daily = read_csv(INFERENCE / "final_model_inference_v1_src_daily_summary.csv")

    lightgbm = leaderboard[leaderboard["model"] == "lightgbm_v1"].iloc[0]
    baseline = leaderboard[
        leaderboard["model"] == "baseline_rolling_mean_28_origin"
    ].iloc[0]
    prophet = leaderboard[leaderboard["model"] == "prophet_total_sales_v1"].iloc[0]
    daily_lightgbm = leaderboard[
        leaderboard["model"] == "lightgbm_v1_daily_aggregate"
    ].iloc[0]

    return {
        "recommended_model": recommendation.iloc[0]["recommended_model"],
        "horizon": int(recommendation.iloc[0]["forecast_horizon_days"]),
        "validation_start": lightgbm["validation_start_date"],
        "validation_end": lightgbm["validation_end_date"],
        "rows": int(lightgbm["rows"]),
        "rmsle": float(lightgbm["rmsle"]),
        "mae": float(lightgbm["mae"]),
        "wape": float(lightgbm["wape"]),
        "bias_pct": float(lightgbm["total_bias_pct"]),
        "baseline_rmsle": float(baseline["rmsle"]),
        "baseline_wape": float(baseline["wape"]),
        "daily_rmsle": float(daily_lightgbm["rmsle"]),
        "prophet_rmsle": float(prophet["rmsle"]),
        "future_rows": int(future_daily["rows"].sum()),
        "future_total": float(future_daily["predicted_sales_sum"].sum()),
    }


def add_title(slide, title: str, subtitle: str | None = None):
    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(8.6), Inches(0.65))
    title_frame = title_box.text_frame
    title_frame.clear()
    title_frame.paragraphs[0].text = title
    title_frame.paragraphs[0].font.size = Pt(30)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(26, 45, 64)

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.58), Inches(1.0), Inches(8.3), Inches(0.35))
        sub_frame = sub_box.text_frame
        sub_frame.clear()
        sub_frame.paragraphs[0].text = subtitle
        sub_frame.paragraphs[0].font.size = Pt(13)
        sub_frame.paragraphs[0].font.color.rgb = RGBColor(80, 88, 96)


def add_bullets(slide, bullets: list[str], left: float, top: float, width: float, height: float):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    frame.clear()

    for index, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.size = Pt(15)
        paragraph.font.color.rgb = RGBColor(32, 42, 52)
        paragraph.space_after = Pt(8)


def add_metric_card(slide, label: str, value: str, left: float, top: float):
    shape = slide.shapes.add_shape(
        1,
        Inches(left),
        Inches(top),
        Inches(1.9),
        Inches(0.82),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(245, 247, 250)
    shape.line.color.rgb = RGBColor(200, 205, 212)
    frame = shape.text_frame
    frame.clear()
    p1 = frame.paragraphs[0]
    p1.text = value
    p1.font.size = Pt(18)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(26, 45, 64)
    p1.alignment = PP_ALIGN.CENTER
    p2 = frame.add_paragraph()
    p2.text = label
    p2.font.size = Pt(8)
    p2.font.color.rgb = RGBColor(80, 88, 96)
    p2.alignment = PP_ALIGN.CENTER


def add_image_if_exists(slide, image_path: Path, left: float, top: float, width: float):
    if image_path.exists():
        slide.shapes.add_picture(
            image_path.as_posix(),
            Inches(left),
            Inches(top),
            width=Inches(width),
        )


def copy_presentation_assets() -> list[Path]:
    PRESENTATION_CHARTS.mkdir(parents=True, exist_ok=True)
    chart_names = [
        "final_model_global_metrics_comparison.png",
        "final_daily_actual_vs_prediction.png",
        "final_worst_stores_by_wape.png",
        "final_worst_families_by_wape.png",
        "final_error_by_date.png",
    ]
    copied = []
    for name in chart_names:
        source = FIGURES / name
        if source.exists():
            target = PRESENTATION_CHARTS / name
            shutil.copy2(source, target)
            copied.append(target)
    return copied


def build_pptx() -> Path:
    metrics = metric_lookup()
    copy_presentation_assets()

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    add_title(slide, TITLE, SUBTITLE)
    add_bullets(
        slide,
        [
            "Goal: forecast 16 days of daily sales at date, store and family grain.",
            "Recommended model: LightGBM global forecasting model.",
            "Serving surface: FastAPI batch prediction endpoint with Docker config.",
        ],
        0.7,
        1.65,
        5.2,
        2.2,
    )
    add_metric_card(slide, "validation rows", f"{metrics['rows']:,}", 6.55, 1.45)
    add_metric_card(slide, "forecast horizon", f"{metrics['horizon']} days", 8.0, 1.45)
    add_metric_card(slide, "future rows", f"{metrics['future_rows']:,}", 6.55, 2.55)
    add_metric_card(slide, "future total sales", f"{metrics['future_total']:,.0f}", 8.0, 2.55)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Business Problem", "Short-term replenishment support")
    add_bullets(
        slide,
        [
            "Forecast expected sales for each store and product family.",
            "Reduce risk of stockouts, excess stock and poor holiday or promotion planning.",
            "Use known future context: calendar, store metadata, promotions, holidays and oil features.",
            "Keep the model interpretable enough for business review and portfolio storytelling.",
        ],
        0.7,
        1.45,
        8.4,
        3.2,
    )

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Model Selection", "LightGBM outperformed the operational baseline")
    add_metric_card(slide, "LightGBM RMSLE", f"{metrics['rmsle']:.4f}", 0.8, 1.35)
    add_metric_card(slide, "Baseline RMSLE", f"{metrics['baseline_rmsle']:.4f}", 3.0, 1.35)
    add_metric_card(slide, "LightGBM WAPE", f"{metrics['wape']:.1%}", 5.2, 1.35)
    add_metric_card(slide, "Baseline WAPE", f"{metrics['baseline_wape']:.1%}", 7.4, 1.35)
    add_image_if_exists(
        slide,
        PRESENTATION_CHARTS / "final_model_global_metrics_comparison.png",
        1.1,
        2.45,
        7.8,
    )

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Daily Forecast Quality", "Aggregate trend check")
    add_bullets(
        slide,
        [
            f"Validation window: {metrics['validation_start']} to {metrics['validation_end']}.",
            f"Daily aggregate LightGBM RMSLE: {metrics['daily_rmsle']:.4f}.",
            f"Daily aggregate Prophet RMSLE: {metrics['prophet_rmsle']:.4f}.",
            f"Total forecast bias: {metrics['bias_pct']:.2%}.",
        ],
        0.65,
        1.35,
        3.2,
        3.6,
    )
    add_image_if_exists(
        slide,
        PRESENTATION_CHARTS / "final_daily_actual_vs_prediction.png",
        4.1,
        1.28,
        5.25,
    )

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Operational Risk Views", "Where planners should monitor forecast error")
    add_image_if_exists(
        slide,
        PRESENTATION_CHARTS / "final_worst_stores_by_wape.png",
        0.55,
        1.25,
        4.25,
    )
    add_image_if_exists(
        slide,
        PRESENTATION_CHARTS / "final_worst_families_by_wape.png",
        5.05,
        1.25,
        4.25,
    )

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Serving And MLOps Readiness", "API-ready model artifacts")
    add_bullets(
        slide,
        [
            "FastAPI exposes liveness, readiness, feature contract, model metadata and batch predict endpoints.",
            "Dockerfile copies only the API, source modules, model files and required serving features.",
            "New tests validate artifact availability, model contract compatibility, API endpoints and input errors.",
            "CI now installs production/dev dependencies, compiles source, runs tests and builds the Docker image.",
        ],
        0.7,
        1.45,
        8.4,
        3.3,
    )

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Next Iteration", "Practical improvements after portfolio handoff")
    add_bullets(
        slide,
        [
            "Build the Power BI report from the generated dashboard/data_exports CSV layer.",
            "Add automated feature generation for raw date, store and family requests if serving non-technical users.",
            "Register the final LightGBM artifact in MLflow and add deployment promotion gates.",
            "Monitor high-risk stores, families, holidays and promotion periods after launch.",
        ],
        0.7,
        1.45,
        8.4,
        3.3,
    )

    output = PRESENTATION_DIR / "store_sales_forecasting_presentation.pptx"
    PRESENTATION_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    return output


def build_pdf() -> Path:
    metrics = metric_lookup()
    output = PRESENTATION_DIR / "store_sales_forecasting_presentation.pdf"
    doc = SimpleDocTemplate(output.as_posix(), pagesize=letter)
    styles = getSampleStyleSheet()
    story = []

    story.append(Paragraph(TITLE, styles["Title"]))
    story.append(Paragraph(SUBTITLE, styles["Normal"]))
    story.append(Spacer(1, 0.2 * inch))

    summary = (
        f"The recommended model is {metrics['recommended_model']} for a "
        f"{metrics['horizon']}-day forecast at date, store and family grain. "
        f"On the validation window {metrics['validation_start']} to "
        f"{metrics['validation_end']}, LightGBM reached RMSLE "
        f"{metrics['rmsle']:.4f}, WAPE {metrics['wape']:.1%}, and total bias "
        f"{metrics['bias_pct']:.2%}."
    )
    story.append(Paragraph(summary, styles["BodyText"]))
    story.append(Spacer(1, 0.2 * inch))

    table_data = [
        ["Metric", "LightGBM", "Reference"],
        ["Granular RMSLE", f"{metrics['rmsle']:.4f}", f"{metrics['baseline_rmsle']:.4f} baseline"],
        ["Granular WAPE", f"{metrics['wape']:.1%}", f"{metrics['baseline_wape']:.1%} baseline"],
        ["Daily RMSLE", f"{metrics['daily_rmsle']:.4f}", f"{metrics['prophet_rmsle']:.4f} Prophet"],
        ["Future forecast rows", f"{metrics['future_rows']:,}", ""],
        ["Future predicted sales", f"{metrics['future_total']:,.0f}", ""],
    ]
    table = Table(table_data, colWidths=[2.0 * inch, 1.8 * inch, 2.2 * inch])
    table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1A2D40")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("GRID", (0, 0), (-1, -1), 0.25, colors.HexColor("#C8CDD4")),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#F5F7FA")]),
            ]
        )
    )
    story.append(table)
    story.append(Spacer(1, 0.25 * inch))

    for chart_name in [
        "final_model_global_metrics_comparison.png",
        "final_daily_actual_vs_prediction.png",
    ]:
        chart_path = FIGURES / chart_name
        if chart_path.exists():
            story.append(Image(chart_path.as_posix(), width=6.4 * inch, height=3.0 * inch))
            story.append(Spacer(1, 0.18 * inch))

    story.append(Paragraph("Operational handoff", styles["Heading2"]))
    story.append(
        Paragraph(
            "Use dashboard/data_exports as the Power BI import layer. The API "
            "expects model-ready feature records and returns non-negative batch "
            "sales forecasts with validation diagnostics.",
            styles["BodyText"],
        )
    )

    doc.build(story)
    return output


def main() -> None:
    exports = build_dashboard_exports()
    pptx = build_pptx()
    pdf = build_pdf()

    print("Generated dashboard exports:")
    for path in exports:
        print(f"- {path.relative_to(ROOT).as_posix()}")
    print(f"Generated presentation: {pptx.relative_to(ROOT).as_posix()}")
    print(f"Generated PDF: {pdf.relative_to(ROOT).as_posix()}")


if __name__ == "__main__":
    main()

